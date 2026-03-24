"""PowerPoint freeform shape creation for polygons, polylines, and paths."""

from typing import Optional, Union

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu

from svg2pptx.parser.shapes import PolygonShape, PolylineShape
from svg2pptx.parser.paths import PathShape
from svg2pptx.config import Config
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import apply_style


def _normalize_emu_points(
    points: list[tuple[int, int]],
    *,
    is_closed: bool,
) -> list[tuple[int, int]]:
    """Drop consecutive duplicates and one repeated closing point before builder IO."""
    normalized: list[tuple[int, int]] = []
    for point in points:
        if not normalized or point != normalized[-1]:
            normalized.append(point)

    if is_closed and len(normalized) > 1 and normalized[0] == normalized[-1]:
        normalized.pop()

    return normalized


def create_freeform(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape, PathShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint freeform shape from polygon, polyline, or path.

    Args:
        shapes: SlideShapes collection to add shape to.
        parsed_shape: Parsed polygon, polyline, or path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created freeform shape or None.
    """
    if isinstance(parsed_shape, PathShape):
        return create_freeform_from_path(
            shapes, parsed_shape, offset_x, offset_y, scale, config=config
        )
    else:
        return create_freeform_from_points(
            shapes, parsed_shape, offset_x, offset_y, scale, config=config
        )


def create_freeform_from_points(
    shapes: SlideShapes,
    parsed_shape: Union[PolygonShape, PolylineShape],
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from polygon or polyline points.

    Args:
        shapes: SlideShapes collection.
        parsed_shape: Parsed polygon or polyline.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    points = parsed_shape.points
    if len(points) < 2:
        return None

    # Apply transform to all points
    transformed_points = parsed_shape.transform.apply_to_points(points)

    # Convert to EMU and apply scale
    emu_points = [
        (px_to_emu(x * scale), px_to_emu(y * scale))
        for x, y in transformed_points
    ]
    # Determine if closed
    is_closed = isinstance(parsed_shape, PolygonShape)
    emu_points = _normalize_emu_points(emu_points, is_closed=is_closed)
    if len(emu_points) < 2:
        return None
    if config is not None:
        config.note_freeform_points(len(emu_points), type(parsed_shape).__name__)

    # Create freeform using FreeformBuilder
    first_x, first_y = emu_points[0]
    builder = shapes.build_freeform(first_x, first_y)
    builder.add_line_segments(emu_points[1:], close=is_closed)

    # Convert to shape
    shape = builder.convert_to_shape(offset_x, offset_y)

    # Apply styling
    apply_style(
        shape,
        parsed_shape.style,
        disable_shadow=config.disable_shadows if config else True,
        config=config,
    )
    if config is not None:
        config.note_shape_created(type(parsed_shape).__name__)

    return shape


def create_freeform_from_path(
    shapes: SlideShapes,
    path: PathShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> Optional[BaseShape]:
    """
    Create a freeform shape from an SVG path.

    Paths with multiple subpaths are combined into a single shape.

    Args:
        shapes: SlideShapes collection.
        path: Parsed path shape.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if not path.subpaths:
        return None

    prepared_subpaths: list[tuple[list[tuple[int, int]], bool]] = []
    for subpath_points, is_closed in path.subpaths:
        if len(subpath_points) < 2:
            continue
        transformed = path.transform.apply_to_points(subpath_points)
        emu_points = [
            (px_to_emu(x * scale), px_to_emu(y * scale))
            for x, y in transformed
        ]
        emu_points = _normalize_emu_points(emu_points, is_closed=is_closed)
        if len(emu_points) < 2:
            continue
        prepared_subpaths.append((emu_points, is_closed))

    if not prepared_subpaths:
        return None

    # Get the first subpath to start the builder
    first_subpath_points, first_closed = prepared_subpaths[0]
    first_x, first_y = first_subpath_points[0]

    # Start the freeform builder
    builder = shapes.build_freeform(first_x, first_y)

    # Add segments for first subpath
    builder.add_line_segments(first_subpath_points[1:], close=first_closed)
    total_points = len(first_subpath_points)

    # Add additional subpaths using move_to
    for emu_points, is_closed in prepared_subpaths[1:]:
        # Move to start of new subpath
        start_x, start_y = emu_points[0]
        builder.move_to(start_x, start_y)

        # Add line segments
        builder.add_line_segments(emu_points[1:], close=is_closed)
        total_points += len(emu_points)

    # Convert to shape
    shape = builder.convert_to_shape(offset_x, offset_y)

    # Apply styling
    apply_style(
        shape,
        path.style,
        disable_shadow=config.disable_shadows if config else True,
        config=config,
    )
    if config is not None:
        config.note_freeform_points(total_points, "PathShape")
        config.note_shape_created("PathShape")

    return shape
