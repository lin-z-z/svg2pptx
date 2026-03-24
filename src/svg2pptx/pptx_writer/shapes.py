"""PowerPoint shape creation utilities."""

import math
from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes, GroupShapes
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from svg2pptx.config import Config
from svg2pptx.parser.styles import FilterSpec, GradientSpec, Style
from svg2pptx.parser.shapes import (
    ParsedShape,
    RectShape,
    CircleShape,
    EllipseShape,
    LineShape,
    PolygonShape,
    PolylineShape,
)
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.geometry.transforms import (
    Transform,
    transform_rect_to_bbox,
    transform_ellipse_to_bbox,
)


def create_shape(
    shapes: SlideShapes,
    parsed_shape: ParsedShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint shape from a parsed SVG shape.

    Args:
        shapes: SlideShapes or GroupShapes collection to add shape to.
        parsed_shape: Parsed shape data.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created shape or None.
    """
    if isinstance(parsed_shape, RectShape):
        return create_rectangle(
            shapes,
            parsed_shape,
            offset_x,
            offset_y,
            scale,
            config=config,
        )
    elif isinstance(parsed_shape, (CircleShape, EllipseShape)):
        return create_oval(
            shapes,
            parsed_shape,
            offset_x,
            offset_y,
            scale,
            config=config,
        )
    elif isinstance(parsed_shape, LineShape):
        return create_line(
            shapes,
            parsed_shape,
            offset_x,
            offset_y,
            scale,
            config=config,
        )
    elif isinstance(parsed_shape, (PolygonShape, PolylineShape)):
        from svg2pptx.pptx_writer.freeform import create_freeform

        return create_freeform(
            shapes,
            parsed_shape,
            offset_x,
            offset_y,
            scale,
            config=config,
        )
    return None


def create_rectangle(
    shapes: SlideShapes,
    rect: RectShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> BaseShape:
    """Create a PowerPoint rectangle shape."""
    x, y, width_px, height_px = transform_rect_to_bbox(
        rect.x,
        rect.y,
        rect.width,
        rect.height,
        rect.transform,
    )

    # Convert to EMU with scale
    left = offset_x + px_to_emu(x * scale)
    top = offset_y + px_to_emu(y * scale)
    width = px_to_emu(width_px * scale)
    height = px_to_emu(height_px * scale)

    # Choose shape type based on corner radius
    if rect.rx > 0 or rect.ry > 0:
        shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        _apply_corner_radius(shape, rect, width_px, height_px, config)
    else:
        shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    apply_style(
        shape,
        rect.style,
        disable_shadow=config.disable_shadows if config else True,
        config=config,
    )
    if config is not None:
        config.note_shape_created("rect")
    return shape


def create_oval(
    shapes: SlideShapes,
    oval: CircleShape | EllipseShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> BaseShape:
    """Create a PowerPoint oval (ellipse/circle) shape."""
    if isinstance(oval, CircleShape):
        left_px, top_px, width_px, height_px = transform_ellipse_to_bbox(
            oval.cx,
            oval.cy,
            oval.r,
            oval.r,
            oval.transform,
        )
    else:
        left_px, top_px, width_px, height_px = transform_ellipse_to_bbox(
            oval.cx,
            oval.cy,
            oval.rx,
            oval.ry,
            oval.transform,
        )

    left = offset_x + px_to_emu(left_px * scale)
    top = offset_y + px_to_emu(top_px * scale)
    width = px_to_emu(width_px * scale)
    height = px_to_emu(height_px * scale)

    shape = shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    apply_style(
        shape,
        oval.style,
        disable_shadow=config.disable_shadows if config else True,
        config=config,
    )
    if config is not None:
        config.note_shape_created("oval")
    return shape


def create_line(
    shapes: SlideShapes,
    line: LineShape,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    config: Optional[Config] = None,
) -> BaseShape:
    """Create a PowerPoint line connector."""
    from pptx.enum.shapes import MSO_CONNECTOR

    # Apply transform
    x1, y1 = line.transform.apply(line.x1, line.y1)
    x2, y2 = line.transform.apply(line.x2, line.y2)

    # Convert to EMU
    start_x = offset_x + px_to_emu(x1 * scale)
    start_y = offset_y + px_to_emu(y1 * scale)
    end_x = offset_x + px_to_emu(x2 * scale)
    end_y = offset_y + px_to_emu(y2 * scale)

    connector = shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )

    _apply_shadow(connector, disable_shadow=config.disable_shadows if config else True)
    _record_style_metadata(line.style, config, source="connector")
    _apply_line_format(connector.line, line.style, config)
    if config is not None:
        config.note_shape_created("line")

    return connector


def apply_style(
    shape: BaseShape,
    style: Style,
    disable_shadow: bool = True,
    config: Optional[Config] = None,
) -> None:
    """
    Apply SVG style to a PowerPoint shape.

    Args:
        shape: PowerPoint shape to style.
        style: Parsed SVG style.
        disable_shadow: Whether to disable shadow on the shape. Defaults to True.
    """
    apply_effects(shape, style, disable_shadow=disable_shadow, config=config)
    _apply_fill_format(shape.fill, style, config)
    _apply_line_format(shape.line, style, config)


def apply_effects(
    shape: BaseShape,
    style: Style,
    disable_shadow: bool = True,
    config: Optional[Config] = None,
) -> None:
    """Apply shadow/filter effects without touching fill or stroke."""
    _apply_shadow(shape, disable_shadow)
    _record_style_metadata(style, config, source="shape")
    _apply_filter_effect(shape, style, config)


def _apply_shadow(shape: BaseShape, disable_shadow: bool = True) -> None:
    """Disable PowerPoint default shadow when requested."""
    if not disable_shadow:
        return

    try:
        shape.shadow.inherit = False
        if hasattr(shape.shadow, "visible"):
            shape.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass


def _apply_filter_effect(
    shape: BaseShape,
    style: Style,
    config: Optional[Config],
) -> None:
    """Approximate supported SVG filters via DrawingML effects."""
    effect = style.filter_effect
    if effect is None:
        return

    sp_pr = getattr(shape._element, "spPr", None)
    if sp_pr is None:
        sp_pr = getattr(shape._element, "grpSpPr", None)
    if sp_pr is None:
        return

    effect_list = sp_pr.find(qn("a:effectLst"))
    if effect_list is None:
        effect_list = OxmlElement("a:effectLst")
        sp_pr.append(effect_list)
    else:
        for child in list(effect_list):
            effect_list.remove(child)

    if effect.kind == "outer_shadow":
        effect_list.append(_build_outer_shadow(effect, style))
        return

    if effect.kind == "glow":
        effect_list.append(_build_glow(effect, style))
        return

    _record_runtime_unsupported(
        config,
        "filter",
        f"url(#{effect.filter_id})",
        "unsupported-filter-chain",
        source="shape",
    )


def _apply_fill_format(fill, style: Style, config: Optional[Config]) -> None:
    """Apply fill color and opacity using one shared code path."""
    if style.fill == "none":
        fill.background()
        return

    if style.fill_gradient is not None:
        fallback_color = _gradient_fallback_color(style.fill_gradient, style, config)
        if _should_fallback_to_solid_fill(style.fill_gradient, style, config):
            if config is not None:
                config.note_gradient_degraded()
            if fallback_color is not None:
                fill.solid()
                fill.fore_color.rgb = parse_hex_color(fallback_color)
                return
        elif _apply_gradient_fill(fill._xPr, style.fill_gradient, style, config):
            return
        elif fallback_color is not None:
            if config is not None:
                config.note_gradient_degraded()
            fill.solid()
            fill.fore_color.rgb = parse_hex_color(fallback_color)
            return

    try:
        color = parse_hex_color(style.fill)
    except ValueError:
        fill.background()
        _record_runtime_unsupported(
            config,
            "fill",
            style.fill,
            "unparsed-color-token",
            source="shape",
        )
        return

    fill.solid()
    fill.fore_color.rgb = color
    _apply_opacity_to_solid_fill(fill._xPr, style.effective_fill_opacity)


def _apply_gradient_fill(
    container,
    gradient: GradientSpec,
    style: Style,
    config: Optional[Config],
) -> bool:
    """Apply a gradient fill to the shape property XML."""
    if gradient.kind not in {"linear", "radial"}:
        _record_runtime_unsupported(
            config,
            "fill",
            f"url(#{gradient.gradient_id})",
            "unsupported-gradient-kind",
            source="shape",
        )
        if config is not None:
            config.note_gradient_degraded()
        return False

    for tag in (
        "a:solidFill",
        "a:gradFill",
        "a:noFill",
        "a:pattFill",
        "a:blipFill",
        "a:grpFill",
    ):
        child = container.find(qn(tag))
        if child is not None:
            container.remove(child)

    grad_fill = OxmlElement("a:gradFill")
    grad_fill.set("flip", "none")
    grad_fill.set("rotWithShape", "1")
    gs_list = OxmlElement("a:gsLst")
    for stop in gradient.stops:
        gs = OxmlElement("a:gs")
        gs.set("pos", str(int(round(max(0.0, min(stop.offset, 1.0)) * 100000))))
        color_node = OxmlElement("a:srgbClr")
        stop_opacity = stop.opacity * style.fill_opacity * style.opacity
        blended_color = _gradient_stop_visible_color(stop.color, stop_opacity, config)
        color_node.set("val", parse_hex_color_to_str(blended_color))
        if blended_color.strip().upper() == stop.color.strip().upper():
            _append_alpha(color_node, stop_opacity)
        gs.append(color_node)
        gs_list.append(gs)
    grad_fill.append(gs_list)

    if gradient.kind == "linear":
        linear = OxmlElement("a:lin")
        linear.set("ang", str(_linear_gradient_angle(gradient)))
        linear.set("scaled", "1")
        grad_fill.append(linear)
        if config is not None:
            config.note_gradient_applied("linear")
    else:
        path = OxmlElement("a:path")
        path.set("path", "circle")
        grad_fill.append(path)
        if config is not None:
            config.note_gradient_applied("radial")
            config.note_gradient_degraded()
        _record_runtime_unsupported(
            config,
            "fill",
            f"url(#{gradient.gradient_id})",
            "radial-gradient-simplified",
            source="shape",
        )
    grad_fill.append(OxmlElement("a:tileRect"))

    _insert_fill_node(container, grad_fill)
    return True


def _should_fallback_to_solid_fill(
    gradient: GradientSpec,
    style: Style,
    config: Optional[Config],
) -> bool:
    """Avoid PowerPoint blue-casting on near-white, low-contrast gradients."""
    if len(gradient.stops) < 2:
        return True

    visible_colors: list[tuple[int, int, int]] = []
    effective_opacities: list[float] = []
    background_rgb = _page_background_rgb(config)
    for stop in gradient.stops:
        rgb = _parse_rgb_triplet(stop.color)
        if rgb is None:
            return False
        effective_opacity = max(
            0.0, min(stop.opacity * style.fill_opacity * style.opacity, 1.0)
        )
        effective_opacities.append(effective_opacity)
        visible_colors.append(
            _blend_rgb_over_background(rgb, effective_opacity, background_rgb)
        )

    if max(effective_opacities, default=1.0) < 0.25:
        return False

    min_channel = min(min(rgb) for rgb in visible_colors)
    max_delta = max(
        abs(a - b)
        for rgb_a in visible_colors
        for rgb_b in visible_colors
        for a, b in zip(rgb_a, rgb_b)
    )
    return min_channel >= 0xE8 and max_delta <= 0x18


def _gradient_fallback_color(
    gradient: GradientSpec,
    style: Style,
    config: Optional[Config] = None,
) -> Optional[str]:
    """Approximate the final visible gradient as one opaque solid color."""
    if not gradient.stops:
        return None

    stops = sorted(gradient.stops, key=lambda stop: stop.offset)
    background_rgb = _page_background_rgb(config)
    accum = [0.0, 0.0, 0.0]
    total_weight = 0.0
    for idx, stop in enumerate(stops):
        rgb = _parse_rgb_triplet(stop.color)
        if rgb is None:
            return None
        effective_opacity = max(
            0.0, min(stop.opacity * style.fill_opacity * style.opacity, 1.0)
        )
        visible_rgb = _blend_rgb_over_background(rgb, effective_opacity, background_rgb)
        weight = _gradient_stop_weight(stops, idx)
        total_weight += weight
        for channel_index, value in enumerate(visible_rgb):
            accum[channel_index] += value * weight

    if total_weight <= 0:
        fallback_rgb = _blend_rgb_over_background(
            _parse_rgb_triplet(stops[0].color) or (255, 255, 255),
            max(0.0, min(stops[0].opacity * style.fill_opacity * style.opacity, 1.0)),
            background_rgb,
        )
    else:
        fallback_rgb = tuple(
            int(round(channel / total_weight)) for channel in accum
        )
    return "#{:02X}{:02X}{:02X}".format(*fallback_rgb)


def _gradient_stop_weight(stops: list, index: int) -> float:
    """Trapezoid-style weight for one stop across a 0..1 gradient span."""
    if len(stops) == 1:
        return 1.0
    prev_offset = 0.0 if index == 0 else max(0.0, min(stops[index - 1].offset, 1.0))
    next_offset = (
        1.0
        if index == len(stops) - 1
        else max(0.0, min(stops[index + 1].offset, 1.0))
    )
    return max((next_offset - prev_offset) / 2.0, 0.0)


def _parse_rgb_triplet(color: str) -> Optional[tuple[int, int, int]]:
    """Parse a strict hex color token into one RGB tuple."""
    normalized = parse_hex_color_to_str(color)
    if len(normalized) != 6:
        return None
    return (
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )


def _blend_rgb_over_white(
    rgb: tuple[int, int, int], opacity: float
) -> tuple[int, int, int]:
    """Composite one RGB color onto a white page background."""
    return _blend_rgb_over_background(rgb, opacity, (255, 255, 255))


def _blend_rgb_over_background(
    rgb: tuple[int, int, int],
    opacity: float,
    background_rgb: tuple[int, int, int],
) -> tuple[int, int, int]:
    """Composite one RGB color onto one resolved page background."""
    clamped = max(0.0, min(opacity, 1.0))
    return tuple(
        int(round(background * (1.0 - clamped) + channel * clamped))
        for channel, background in zip(rgb, background_rgb)
    )


def _gradient_stop_visible_color(
    color: str,
    opacity: float,
    config: Optional[Config],
) -> str:
    """Convert one potentially translucent stop into the visible page color."""
    if opacity >= 0.9999:
        return color
    rgb = _parse_rgb_triplet(color)
    if rgb is None:
        return color
    blended = _blend_rgb_over_background(rgb, opacity, _page_background_rgb(config))
    return "#{:02X}{:02X}{:02X}".format(*blended)


def _page_background_rgb(config: Optional[Config]) -> tuple[int, int, int]:
    """Resolve the current page background for color precompositing."""
    token = "#FFFFFF" if config is None else getattr(config, "page_background", "#FFFFFF")
    rgb = _parse_rgb_triplet(token)
    return rgb or (255, 255, 255)


def _build_outer_shadow(effect: FilterSpec, style: Style):
    """Build an outer shadow effect from a parsed filter spec."""
    shadow = OxmlElement("a:outerShdw")
    shadow.set("blurRad", str(int(round(px_to_emu(max(effect.blur * 1.5, 0.0))))))
    distance = math.hypot(effect.dx, effect.dy)
    shadow.set("dist", str(int(round(px_to_emu(distance)))))
    shadow.set("dir", str(_vector_angle(effect.dx, effect.dy)))
    shadow.set("rotWithShape", "0")

    color_node = OxmlElement("a:srgbClr")
    color_node.set("val", _effect_color_hex(effect, style))
    _append_alpha(color_node, max(0.0, min(effect.opacity * style.opacity, 1.0)))
    shadow.append(color_node)
    return shadow


def _build_glow(effect: FilterSpec, style: Style):
    """Build a glow effect from a parsed filter spec."""
    glow = OxmlElement("a:glow")
    glow.set("rad", str(int(round(px_to_emu(max(effect.blur * 2.0, 0.0))))))

    color_node = OxmlElement("a:srgbClr")
    color_node.set("val", _effect_color_hex(effect, style))
    base_opacity = style.effective_fill_opacity
    if style.fill == "none":
        base_opacity = style.effective_stroke_opacity
    _append_alpha(color_node, max(0.0, min(effect.opacity * base_opacity, 1.0)))
    glow.append(color_node)
    return glow


def _effect_color_hex(effect: FilterSpec, style: Style) -> str:
    """Resolve the final effect color in RRGGBB form."""
    if effect.use_source_color:
        if style.fill not in ("", "none"):
            return parse_hex_color_to_str(style.fill)
        if style.stroke not in ("", "none"):
            return parse_hex_color_to_str(style.stroke)
        return "000000"
    if effect.color and effect.color != "none":
        return parse_hex_color_to_str(effect.color)
    return "000000"


def _apply_line_format(line_format, style: Style, config: Optional[Config]) -> None:
    """Apply stroke color, width, and opacity using one shared code path."""
    if style.stroke == "none":
        line_format.fill.background()
        return

    try:
        color = parse_hex_color(style.stroke)
    except ValueError:
        line_format.fill.background()
        _record_runtime_unsupported(
            config,
            "stroke",
            style.stroke,
            "unparsed-color-token",
            source="shape",
        )
        return

    line_format.color.rgb = color
    line_format.width = Emu(px_to_emu(style.stroke_width))
    if line_format._ln is not None:
        _apply_opacity_to_solid_fill(line_format._ln, style.effective_stroke_opacity)


def _apply_opacity_to_solid_fill(container, opacity: float) -> None:
    """Write SVG opacity to DrawingML alpha under the solid fill color node."""
    solid_fill = container.find(qn("a:solidFill"))
    if solid_fill is None or len(solid_fill) == 0:
        return

    color_choice = solid_fill[0]
    _append_alpha(color_choice, opacity)


def _append_alpha(color_choice, opacity: float) -> None:
    """Replace or append DrawingML alpha on a color node."""
    for alpha in color_choice.findall(qn("a:alpha")):
        color_choice.remove(alpha)

    clamped = max(0.0, min(opacity, 1.0))
    if clamped >= 0.9999:
        return

    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(int(round(clamped * 100000))))
    color_choice.append(alpha)


def _linear_gradient_angle(gradient: GradientSpec) -> int:
    """Map SVG linear gradient vector to DrawingML ang units (1/60000 deg)."""
    dx = gradient.x2 - gradient.x1
    dy = gradient.y2 - gradient.y1
    if abs(dx) < 1e-9 and abs(dy) < 1e-9:
        return 0
    degrees = math.degrees(math.atan2(dy, dx)) % 360.0
    return int(round(degrees * 60000))


def _insert_fill_node(container, fill_node) -> None:
    """Insert one fill node before line/effect children to satisfy PowerPoint."""
    line_node = container.find(qn("a:ln"))
    if line_node is not None:
        line_index = list(container).index(line_node)
        container.insert(line_index, fill_node)
        return

    effect_node = container.find(qn("a:effectLst"))
    if effect_node is not None:
        effect_index = list(container).index(effect_node)
        container.insert(effect_index, fill_node)
        return

    container.append(fill_node)


def _vector_angle(dx: float, dy: float) -> int:
    """Map an XY vector to DrawingML angle units."""
    if abs(dx) < 1e-9 and abs(dy) < 1e-9:
        return 0
    degrees = math.degrees(math.atan2(dy, dx)) % 360.0
    return int(round(degrees * 60000))


def parse_hex_color_to_str(hex_color: str) -> str:
    """Normalize a hex color string to upper-case RRGGBB."""
    return hex_color.strip().lstrip("#").upper()


def _apply_corner_radius(
    shape: BaseShape,
    rect: RectShape,
    width_px: float,
    height_px: float,
    config: Optional[Config],
) -> None:
    """Map SVG rx/ry to the single rounded-rectangle adjustment that PPT supports."""
    if not getattr(shape, "adjustments", None):
        return

    radii = [radius for radius in (rect.rx, rect.ry) if radius > 0]
    if not radii:
        return

    min_dimension = min(width_px, height_px)
    if min_dimension <= 0:
        return

    radius = min(radii)
    shape.adjustments[0] = max(0.0, min(radius / min_dimension, 0.5))

    if rect.rx > 0 and rect.ry > 0 and abs(rect.rx - rect.ry) > 0.01:
        _record_runtime_unsupported(
            config,
            "radius",
            f"rx={rect.rx},ry={rect.ry}",
            "non-uniform-radius-fallback",
            source="shape",
        )


def _record_style_metadata(
    style: Style,
    config: Optional[Config],
    source: str,
) -> None:
    """Forward parser-stage unsupported style metadata into runtime diagnostics."""
    if config is None:
        return

    for item in style.unsupported_styles:
        config.record_unsupported_style(
            item["property"],
            item["value"],
            item["reason"],
            source=source,
        )


def _record_runtime_unsupported(
    config: Optional[Config],
    property_name: str,
    value: str,
    reason: str,
    source: str,
) -> None:
    """Record a downgraded runtime style item when config is available."""
    if config is None:
        return
    config.record_unsupported_style(property_name, value, reason, source=source)


def parse_hex_color(hex_color: str) -> RGBColor:
    """
    Parse a hex color string to RGBColor.

    Args:
        hex_color: Color in format "#RRGGBB" or "#RGB".

    Returns:
        RGBColor object.

    Raises:
        ValueError: If color format is invalid.
    """
    if not hex_color or hex_color == "none":
        raise ValueError("Invalid color: none")

    color = hex_color.strip().lstrip("#")

    if len(color) == 3:
        color = "".join(c * 2 for c in color)

    if len(color) != 6:
        raise ValueError(f"Invalid hex color: {hex_color}")

    try:
        r = int(color[0:2], 16)
        g = int(color[2:4], 16)
        b = int(color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        raise ValueError(f"Invalid hex color: {hex_color}")
