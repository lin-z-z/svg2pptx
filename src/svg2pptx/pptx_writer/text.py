"""PowerPoint text box creation from SVG text elements."""

from typing import Optional

from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from svg2pptx.parser.svg_parser import TextElement, TextSpan
from svg2pptx.geometry.units import px_to_emu
from svg2pptx.pptx_writer.shapes import parse_hex_color


def create_text(
    shapes: SlideShapes,
    text_element: TextElement,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
) -> Optional[BaseShape]:
    """
    Create a PowerPoint text box from an SVG text element.

    Args:
        shapes: SlideShapes collection to add text box to.
        text_element: Parsed SVG text element.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.

    Returns:
        Created text box shape or None.
    """
    if not text_element.text:
        return None

    paragraphs = _group_spans_into_paragraphs(text_element)

    # Apply transform to position
    x, y = text_element.transform.apply(text_element.x, text_element.y)

    estimated_width, estimated_height, first_font_size = (
        _estimate_text_box_metrics(paragraphs, scale)
    )

    # Convert position to EMU
    # In SVG, text y-coordinate is the baseline position
    # We need to adjust based on the text-anchor for horizontal positioning
    text_anchor = text_element.style.text_anchor
    
    # Calculate left position based on text-anchor
    if text_anchor == "middle":
        # Text is centered at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width / 2)
    elif text_anchor == "end":
        # Text ends at x
        left = offset_x + px_to_emu(x * scale) - px_to_emu(estimated_width)
    else:
        # Default: text starts at x (text-anchor="start")
        left = offset_x + px_to_emu(x * scale)

    # For y position: SVG y is the baseline, so we move up by approximately
    # the font ascent (roughly 80% of font size for most fonts)
    baseline_offset = first_font_size * 0.85
    top = offset_y + px_to_emu(y * scale) - px_to_emu(baseline_offset)
    
    width = px_to_emu(estimated_width)
    height = px_to_emu(estimated_height)

    # Create text box
    text_box = shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = False
    
    # Remove margins/padding for more accurate positioning
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    text_frame.clear()
    anchor_map = {
        "start": PP_ALIGN.LEFT,
        "middle": PP_ALIGN.CENTER,
        "end": PP_ALIGN.RIGHT,
    }
    alignment = anchor_map.get(text_anchor, PP_ALIGN.LEFT)

    for paragraph_index, spans in enumerate(paragraphs):
        paragraph = (
            text_frame.paragraphs[0]
            if paragraph_index == 0
            else text_frame.add_paragraph()
        )
        paragraph.alignment = alignment

        for span in spans:
            run = paragraph.add_run()
            run.text = span.text
            _apply_run_style(run, span.style, scale)

    # Disable shadow on text box
    try:
        text_box.shadow.inherit = False
        if hasattr(text_box.shadow, 'visible'):
            text_box.shadow.visible = False
    except (AttributeError, NotImplementedError):
        pass

    return text_box


def _group_spans_into_paragraphs(text_element: TextElement) -> list[list[TextSpan]]:
    """Group spans into paragraph-sized chunks using y/dy semantics."""
    spans = text_element.spans or [
        TextSpan(
            text=text_element.text,
            style=text_element.style,
            x=text_element.x,
            y=text_element.y,
        )
    ]

    paragraphs: list[list[TextSpan]] = []
    current_paragraph: list[TextSpan] = []
    previous_span: Optional[TextSpan] = None

    for span in spans:
        if previous_span is not None and _starts_new_paragraph(span, previous_span):
            paragraphs.append(current_paragraph)
            current_paragraph = []
        current_paragraph.append(span)
        previous_span = span

    if current_paragraph:
        paragraphs.append(current_paragraph)

    return paragraphs


def _starts_new_paragraph(current: TextSpan, previous: TextSpan) -> bool:
    """Treat dy/y breaks as paragraph boundaries, keep inline tspans as runs."""
    if current.dy != 0:
        return True
    if current.y is not None and previous.y is not None and current.y != previous.y:
        return True
    return False


def _estimate_text_box_metrics(
    paragraphs: list[list[TextSpan]],
    scale: float,
) -> tuple[float, float, float]:
    """Estimate textbox width, height, and first-line font size in pixels."""
    if not paragraphs:
        return (0.0, 0.0, 12.0)

    paragraph_heights: list[float] = []
    paragraph_widths: list[float] = []

    for spans in paragraphs:
        max_font_size = max(span.style.font_size for span in spans) * scale
        paragraph_heights.append(max_font_size * 1.4)
        paragraph_widths.append(
            sum(_estimate_span_width(span, scale) for span in spans) + max_font_size
        )

    first_font_size = max(span.style.font_size for span in paragraphs[0]) * scale
    total_height = paragraph_heights[0]

    for paragraph_index in range(1, len(paragraphs)):
        dy_gap = max(paragraphs[paragraph_index][0].dy * scale, 0.0)
        total_height += max(dy_gap, paragraph_heights[paragraph_index])

    return (max(paragraph_widths), total_height, first_font_size)


def _estimate_span_width(span: TextSpan, scale: float) -> float:
    """Estimate span width in pixels for textbox sizing."""
    font_size_px = span.style.font_size * scale
    char_count = len(span.text)
    base_width = char_count * font_size_px * 0.6
    letter_spacing = max(char_count - 1, 0) * span.style.letter_spacing * scale
    dx_padding = max(span.dx * scale, 0.0)
    return base_width + letter_spacing + dx_padding


def _apply_run_style(run, style, scale: float) -> None:
    """Apply text styling to a run."""
    font = run.font
    font.name = style.font_family
    font.size = Pt(style.font_size * scale)

    if style.font_weight in ("bold", "700", "800", "900"):
        font.bold = True

    if style.fill != "none":
        try:
            color = parse_hex_color(style.fill)
            font.color.rgb = color
        except ValueError:
            pass

