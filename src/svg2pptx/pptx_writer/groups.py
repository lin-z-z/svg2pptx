"""PowerPoint group shape creation."""

from dataclasses import replace
from typing import Optional

from pptx.shapes.group import GroupShape
from pptx.shapes.shapetree import SlideShapes

from svg2pptx.parser.svg_parser import GroupElement, TextElement
from svg2pptx.parser.shapes import ParsedShape
from svg2pptx.parser.paths import PathShape
from svg2pptx.pptx_writer.shapes import apply_effects, create_shape
from svg2pptx.pptx_writer.freeform import create_freeform
from svg2pptx.config import Config


def create_group(
    shapes: SlideShapes,
    group: GroupElement,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    flatten: bool = False,
    config: Optional[Config] = None,
) -> Optional[GroupShape]:
    """
    Create a PowerPoint group from an SVG group element.

    Args:
        shapes: SlideShapes collection to add group to.
        group: Parsed SVG group element.
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.
        flatten: If True, don't create groups, just add shapes directly.
        config: Optional configuration settings.

    Returns:
        Created GroupShape or None.
    """
    if not group.children:
        return None

    if flatten:
        # Add shapes directly without grouping.
        pending_filter = group.style.filter_effect
        for child in group.children:
            if pending_filter is not None:
                child, applied = _attach_group_filter(child, pending_filter)
                if applied:
                    pending_filter = None
            add_element_to_shapes(
                shapes, child, offset_x, offset_y, scale, flatten, config
            )
        return None

    # Create a new group
    group_shape = shapes.add_group_shape()
    if config is not None:
        config.note_shape_created("group")
    group_shapes = group_shape.shapes

    # Add children to the group
    for child in group.children:
        add_element_to_shapes(
            group_shapes, child, offset_x, offset_y, scale, flatten, config
        )

    if group.style.filter_effect is not None:
        apply_effects(
            group_shape,
            group.style,
            disable_shadow=config.disable_shadows if config else True,
            config=config,
        )

    return group_shape


def add_element_to_shapes(
    shapes: SlideShapes,
    element,
    offset_x: int = 0,
    offset_y: int = 0,
    scale: float = 1.0,
    flatten: bool = False,
    config: Optional[Config] = None,
) -> None:
    """
    Add a parsed element to a shapes collection.

    Handles shapes, paths, groups, and text elements.

    Args:
        shapes: SlideShapes or GroupShapes collection.
        element: Parsed element (ParsedShape, PathShape, GroupElement, or TextElement).
        offset_x: X offset in EMU.
        offset_y: Y offset in EMU.
        scale: Scale factor.
        flatten: If True, don't create groups.
        config: Optional configuration settings.
    """
    from svg2pptx.parser.svg_parser import GroupElement, TextElement
    from svg2pptx.pptx_writer.text import create_text

    # Default config if not provided
    if config is None:
        config = Config()

    if isinstance(element, GroupElement):
        create_group(shapes, element, offset_x, offset_y, scale, flatten, config)
    elif isinstance(element, PathShape):
        if config.convert_shapes:
            create_freeform(
                shapes,
                element,
                offset_x,
                offset_y,
                scale,
                config=config,
            )
    elif isinstance(element, TextElement):
        if config.convert_text:
            create_text(shapes, element, offset_x, offset_y, scale, config=config)
    elif isinstance(element, ParsedShape):
        if config.convert_shapes:
            create_shape(
                shapes,
                element,
                offset_x,
                offset_y,
                scale,
                config=config,
            )


def _attach_group_filter(element, filter_effect) -> tuple[object, bool]:
    """Attach one group-level filter to the primary drawable child when flattening."""
    if isinstance(element, ParsedShape):
        return replace(
            element,
            style=replace(element.style, filter_effect=filter_effect),
        ), True
    if isinstance(element, PathShape):
        return replace(
            element,
            style=replace(element.style, filter_effect=filter_effect),
        ), True
    if isinstance(element, TextElement):
        return element, False
    if isinstance(element, GroupElement):
        new_children = []
        applied = False
        for child in element.children:
            if not applied:
                child, applied = _attach_group_filter(child, filter_effect)
            new_children.append(child)
        if applied:
            return replace(element, children=new_children), True
        return element, False
    return element, False
