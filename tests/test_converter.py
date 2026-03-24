"""Tests for the main SVG to PPTX converter."""

from xml.etree import ElementTree as ET
import pytest
from pathlib import Path
import tempfile
import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from svg2pptx import svg_to_pptx, SVGConverter, Config
from svg2pptx.geometry.units import px_to_emu, px_to_pt
from svg2pptx.parser.paths import parse_path
from svg2pptx.pptx_writer.shapes import parse_hex_color
from svg2pptx.pptx_writer.text import _estimate_span_width
from svg2pptx.parser.svg_parser import TextSpan
from svg2pptx.parser.styles import Style


FIXTURES_DIR = Path(__file__).parent / "fixtures"


class TestSVGConverter:
    """Tests for SVGConverter class."""

    def test_convert_simple_svg_string(self):
        """Test converting a simple SVG string."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="10" y="10" width="80" height="80" fill="red"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert prs is not None
        assert len(prs.slides) == 1
        # Should have at least one shape
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_circle(self):
        """Test converting a circle."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <circle cx="50" cy="50" r="40" fill="blue"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        assert len(prs.slides[0].shapes) >= 1

    def test_convert_with_config(self):
        """Test conversion with custom config."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50"/>
        </svg>'''
        
        config = Config(scale=2.0)
        converter = SVGConverter(config=config)
        prs = converter.convert_string(svg)
        
        assert prs is not None

    def test_convert_multiple_shapes(self):
        """Test converting multiple shapes."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">
            <rect x="10" y="10" width="50" height="50" fill="red"/>
            <circle cx="120" cy="35" r="25" fill="blue"/>
            <ellipse cx="170" cy="60" rx="20" ry="30" fill="green"/>
        </svg>'''
        
        converter = SVGConverter()
        prs = converter.convert_string(svg)
        
        # Should have 3 shapes
        assert len(prs.slides[0].shapes) >= 3


class TestConvertFile:
    """Tests for file-based conversion."""

    def test_convert_basic_shapes(self):
        """Test converting the basic shapes fixture."""
        svg_path = FIXTURES_DIR / "basic_shapes.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            # Verify the output file exists and is valid
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides) == 1
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_path_icon(self):
        """Test converting a path-based icon."""
        svg_path = FIXTURES_DIR / "path_icon.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            svg_to_pptx(str(svg_path), pptx_path)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            assert len(prs.slides[0].shapes) > 0
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_convert_grouped_svg(self):
        """Test converting SVG with groups."""
        svg_path = FIXTURES_DIR / "grouped.svg"
        if not svg_path.exists():
            pytest.skip("Fixture file not found")
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            config = Config(preserve_groups=True)
            svg_to_pptx(str(svg_path), pptx_path, config=config)
            
            assert os.path.exists(pptx_path)
            prs = Presentation(pptx_path)
            top_shapes = prs.slides[0].shapes
            assert len(top_shapes) == 2
            assert all(
                shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in top_shapes
            )
            assert len(top_shapes[0].shapes) == 2
            assert top_shapes[0].shapes[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            assert top_shapes[0].shapes[1].shape_type == MSO_SHAPE_TYPE.GROUP
        finally:
            if os.path.exists(pptx_path):
                os.unlink(pptx_path)

    def test_grouped_writer_preserves_shape_and_text_order(self):
        """Grouped writer should keep the original child order for z-order fidelity."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="160" height="120">
            <g transform="translate(10, 10)">
                <rect x="0" y="0" width="120" height="60" fill="#0ea5e9"/>
                <text x="16" y="36" font-size="20" fill="#ffffff">Front label</text>
            </g>
        </svg>"""

        prs = SVGConverter(Config(preserve_groups=True)).convert_string(svg)

        group_shape = prs.slides[0].shapes[0]
        assert group_shape.shape_type == MSO_SHAPE_TYPE.GROUP
        assert [child.shape_type for child in group_shape.shapes] == [
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.TEXT_BOX,
        ]

    def test_grouped_and_flattened_share_nested_group_bbox(self):
        """Nested group geometry should stay consistent across both writers."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="200">
            <g transform="translate(10, 20)">
                <g transform="scale(2)">
                    <rect x="5" y="7" width="10" height="15" fill="#ff0000"/>
                </g>
            </g>
        </svg>"""

        flattened = SVGConverter(Config(flatten_groups=True)).convert_string(svg)
        grouped = SVGConverter(
            Config(preserve_groups=True, flatten_groups=False)
        ).convert_string(svg)

        flat_shape = flattened.slides[0].shapes[0]
        grouped_shape = grouped.slides[0].shapes[0].shapes[0].shapes[0]

        expected_left = Emu(px_to_emu(20))
        expected_top = Emu(px_to_emu(34))
        expected_width = Emu(px_to_emu(20))
        expected_height = Emu(px_to_emu(30))

        assert flat_shape.left == expected_left
        assert flat_shape.top == expected_top
        assert flat_shape.width == expected_width
        assert flat_shape.height == expected_height

        assert grouped_shape.left == expected_left
        assert grouped_shape.top == expected_top
        assert grouped_shape.width == expected_width
        assert grouped_shape.height == expected_height

    def test_viewbox_mapping_scales_x_and_y_consistently(self):
        """viewBox mapping should use one document transform for all shape types."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100" viewBox="0 0 100 100">
            <rect x="10" y="20" width="30" height="40" fill="#ff0000"/>
            <circle cx="50" cy="50" r="10" fill="#00ff00"/>
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        rect = prs.slides[0].shapes[0]
        circle = prs.slides[0].shapes[1]

        assert rect.left == Emu(px_to_emu(20))
        assert rect.top == Emu(px_to_emu(20))
        assert rect.width == Emu(px_to_emu(60))
        assert rect.height == Emu(px_to_emu(40))

        assert circle.left == Emu(px_to_emu(80))
        assert circle.top == Emu(px_to_emu(40))
        assert circle.width == Emu(px_to_emu(40))
        assert circle.height == Emu(px_to_emu(20))

    def test_text_writer_preserves_inline_tspan_runs(self):
        """Inline tspans should become editable runs in one paragraph."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <text x="20" y="40" font-family="Noto Sans SC" font-size="32" fill="#0F172A">
                5-10<tspan font-size="48" fill="#0E5A8A" dx="8">x</tspan>
            </text>
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        text_shape = prs.slides[0].shapes[0]
        paragraph = text_shape.text_frame.paragraphs[0]

        assert len(text_shape.text_frame.paragraphs) == 1
        assert [run.text for run in paragraph.runs] == ["5-10", "x"]
        assert paragraph.runs[0].font.size.pt == pytest.approx(px_to_pt(32))
        assert paragraph.runs[1].font.size.pt == pytest.approx(px_to_pt(48))
        assert paragraph.runs[1].font.color.rgb == parse_hex_color("#0E5A8A")

    def test_text_writer_turns_multiline_tspan_into_paragraphs(self):
        """dy-driven tspans should stay inside one textbox as multiple paragraphs."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="400" height="180">
            <text x="20" y="40" font-family="Noto Sans SC" font-size="18" fill="#64748B">
                <tspan x="20" dy="0">第一行</tspan>
                <tspan x="20" dy="28">第二行</tspan>
                <tspan x="20" dy="28">第三行</tspan>
            </text>
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        text_shape = prs.slides[0].shapes[0]
        paragraphs = text_shape.text_frame.paragraphs

        assert len(paragraphs) == 3
        assert [paragraph.text for paragraph in paragraphs] == [
            "第一行",
            "第二行",
            "第三行",
        ]

    def test_text_writer_handles_oceanppt_fixture_tspans(self):
        """Real OceanPPT text pages should produce editable multi-run text boxes."""
        svg_path = FIXTURES_DIR / "oceanppt" / "baseline_5" / "slide_006.svg"
        converter = SVGConverter()
        presentation = converter.convert_string(svg_path.read_text(encoding="utf-8"))

        text_shapes = [
            shape for shape in presentation.slides[0].shapes if hasattr(shape, "text_frame")
        ]
        assert any(
            len(shape.text_frame.paragraphs) >= 3
            for shape in text_shapes
        )
        assert any(
            any(len(paragraph.runs) >= 2 for paragraph in shape.text_frame.paragraphs)
            for shape in text_shapes
        )

    def test_text_width_estimator_treats_cjk_as_wider_than_latin(self):
        """Chinese characters should reserve more width than latin labels."""
        cjk_span = TextSpan(text="感谢聆听", style=Style(font_size=82))
        latin_span = TextSpan(text="TEST", style=Style(font_size=82))

        assert _estimate_span_width(cjk_span, 1.0) > _estimate_span_width(
            latin_span, 1.0
        )

    def test_shape_writer_maps_opacity_and_corner_radius(self):
        """Basic fill/stroke opacity and rounded radius should be written into PPT XML."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <rect
                x="20"
                y="10"
                width="100"
                height="50"
                rx="10"
                ry="20"
                fill="rgba(255, 0, 0, 0.5)"
                stroke="#00ff00"
                stroke-opacity="0.25"
                stroke-width="4"
            />
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        shape = prs.slides[0].shapes[0]

        fill_color = shape.fill._xPr.find(qn("a:solidFill"))[0]
        line_color = shape.line._ln.find(qn("a:solidFill"))[0]

        assert fill_color.find(qn("a:alpha")).get("val") == "50000"
        assert line_color.find(qn("a:alpha")).get("val") == "25000"
        assert shape.adjustments[0] == pytest.approx(0.2, abs=1e-4)

    def test_shape_writer_maps_linear_gradient_fill(self):
        """linearGradient defs should become DrawingML gradFill rather than a solid fallback."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <defs>
                <linearGradient id="heroGrad" x1="0%" y1="0%" x2="100%" y2="0%">
                    <stop offset="0%" stop-color="#0E5A8A" stop-opacity="0.2" />
                    <stop offset="100%" stop-color="#1E88E5" />
                </linearGradient>
            </defs>
            <rect x="20" y="10" width="100" height="50" fill="url(#heroGrad)" />
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        shape = prs.slides[0].shapes[0]

        grad_fill = shape.fill._xPr.find(qn("a:gradFill"))
        assert grad_fill is not None

        stops = grad_fill.find(qn("a:gsLst"))
        assert len(stops) == 2
        assert stops[0].get("pos") == "0"
        assert stops[0][0].get("val") == "CFDEE8"
        assert stops[0][0].find(qn("a:alpha")) is None
        assert stops[1][0].get("val") == "1E88E5"

        linear = grad_fill.find(qn("a:lin"))
        assert linear is not None
        assert linear.get("ang") == "0"
        assert grad_fill.find(qn("a:tileRect")) is not None
        sppr_children = [child.tag for child in shape.fill._xPr]
        assert sppr_children.index(qn("a:gradFill")) < sppr_children.index(qn("a:ln"))

    def test_shape_writer_degrades_near_white_gradient_to_solid_fill(self):
        """Low-contrast near-white gradients should fall back to solid fill."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <defs>
                <linearGradient id="bgGrad" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" stop-color="#FFFFFF" />
                    <stop offset="100%" stop-color="#F0F7FF" />
                </linearGradient>
            </defs>
            <rect x="20" y="10" width="100" height="50" fill="url(#bgGrad)" />
        </svg>"""

        converter = SVGConverter()
        prs = converter.convert_string(svg)
        shape = prs.slides[0].shapes[0]

        assert shape.fill._xPr.find(qn("a:gradFill")) is None
        solid_fill = shape.fill._xPr.find(qn("a:solidFill"))
        assert solid_fill is not None
        assert solid_fill[0].get("val") == "F8FBFF"
        assert converter.config.gradient_stats["degraded"] == 1

    def test_shape_writer_preblends_translucent_gradient_stops_against_page_background(
        self,
    ):
        """Translucent gradients should keep geometry but use visible stop colors."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <rect width="200" height="120" fill="#F8FCFF" />
            <defs>
                <linearGradient id="heroGrad" x1="0%" y1="0%" x2="100%" y2="100%">
                    <stop offset="0%" stop-color="#0E5A8A" stop-opacity="0.05" />
                    <stop offset="100%" stop-color="#0E5A8A" stop-opacity="0.15" />
                </linearGradient>
            </defs>
            <rect x="20" y="10" width="100" height="50" fill="url(#heroGrad)" opacity="0.4" />
        </svg>"""

        converter = SVGConverter()
        prs = converter.convert_string(svg)
        shape = prs.slides[0].shapes[1]

        grad_fill = shape.fill._xPr.find(qn("a:gradFill"))
        assert grad_fill is not None
        stops = grad_fill.find(qn("a:gsLst"))
        assert [stop[0].get("val") for stop in stops] == ["F3F9FD", "EAF2F8"]
        assert stops[0][0].find(qn("a:alpha")) is None
        assert stops[1][0].find(qn("a:alpha")) is None
        assert converter.config.gradient_stats["linear_applied"] == 1

    def test_shape_writer_maps_drop_shadow_filter_to_outer_shadow(self):
        """Supported shadow filters should become DrawingML outerShdw effects."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <defs>
                <filter id="cardShadow">
                    <feDropShadow dx="0" dy="8" stdDeviation="12" flood-color="#0E5A8A" flood-opacity="0.08" />
                </filter>
            </defs>
            <g filter="url(#cardShadow)">
                <rect x="20" y="10" width="100" height="50" fill="#ffffff" />
                <text x="30" y="40" font-size="18" fill="#0F172A">Title</text>
            </g>
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        rect = prs.slides[0].shapes[0]
        text = prs.slides[0].shapes[1]

        effect_list = rect._element.spPr.find(qn("a:effectLst"))
        shadow = effect_list.find(qn("a:outerShdw"))
        assert shadow is not None
        assert shadow.get("dist") == str(px_to_emu(8))
        assert shadow.get("dir") == str(90 * 60000)
        assert shadow.find(qn("a:srgbClr")).get("val") == "0E5A8A"
        assert shadow.find(qn("a:srgbClr")).find(qn("a:alpha")).get("val") == "8000"
        assert text._element.spPr.find(qn("a:effectLst")).find(qn("a:outerShdw")) is None

    def test_shape_writer_maps_blur_composite_filter_to_glow(self):
        """Blur plus composite should become a glow effect instead of staying silent."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <defs>
                <filter id="glow">
                    <feGaussianBlur stdDeviation="6" result="blur" />
                    <feComposite in="SourceGraphic" in2="blur" operator="over" />
                </filter>
            </defs>
            <rect x="20" y="10" width="100" height="50" fill="#0E5A8A" filter="url(#glow)" />
        </svg>"""

        prs = SVGConverter().convert_string(svg)
        rect = prs.slides[0].shapes[0]

        effect_list = rect._element.spPr.find(qn("a:effectLst"))
        glow = effect_list.find(qn("a:glow"))
        assert glow is not None
        assert glow.get("rad") == str(px_to_emu(12))
        assert glow.find(qn("a:srgbClr")).get("val") == "0E5A8A"

    def test_converter_records_shape_count_overflow_warning(self):
        """Shape inflation should emit a structured warning instead of staying silent."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="200" height="120">
            <rect x="10" y="10" width="20" height="20" fill="#ff0000" />
            <rect x="40" y="10" width="20" height="20" fill="#00ff00" />
            <rect x="70" y="10" width="20" height="20" fill="#0000ff" />
        </svg>"""

        converter = SVGConverter(Config(max_shapes_per_page=2))
        converter.convert_string(svg)

        assert converter.config.render_metrics["shape_count"] == 3
        assert any(
            warning["code"] == "shape-count-overflow"
            for warning in converter.config.render_warnings
        )

    def test_converter_records_freeform_point_overflow_warning(self):
        """Large freeform point sets should surface page and per-shape warnings."""
        points = " ".join(f"{idx},{idx % 7}" for idx in range(80))
        svg = f'''<svg xmlns="http://www.w3.org/2000/svg" width="300" height="120">
            <polyline points="{points}" fill="none" stroke="#0E5A8A" />
        </svg>'''

        converter = SVGConverter(
            Config(max_points_per_freeform=40, max_freeform_points_per_page=40)
        )
        converter.convert_string(svg)

        assert converter.config.render_metrics["freeform_points"] == 80
        codes = {warning["code"] for warning in converter.config.render_warnings}
        assert "freeform-points-per-shape-overflow" in codes
        assert "freeform-points-per-page-overflow" in codes

    def test_closed_path_drops_duplicate_closing_point_before_freeform_write(self):
        """Closed path output should not emit one extra segment back to the start point."""
        svg = """<svg xmlns="http://www.w3.org/2000/svg" width="120" height="120">
            <path d="M 10 10 C 70 0, 110 40, 110 80 L 20 100 Z" fill="#0E5A8A" />
        </svg>"""
        root = ET.fromstring(svg)
        parsed = parse_path(root.find("{http://www.w3.org/2000/svg}path"))
        assert parsed is not None
        raw_points, is_closed = parsed.subpaths[0]
        assert is_closed is True
        assert raw_points[0] == raw_points[-1]

        converter = SVGConverter(Config(curve_tolerance=1.0))
        converter.convert_string(svg)

        assert converter.config.render_metrics["freeform_points"] == len(raw_points) - 1

    def test_centered_cjk_title_gets_non_trivial_textbox_width(self):
        """Real centered Chinese titles should not collapse to latin-style width."""
        svg_path = FIXTURES_DIR / "oceanppt" / "full_15" / "slide_015.svg"
        converter = SVGConverter()
        presentation = converter.convert_string(svg_path.read_text(encoding="utf-8"))

        title_shape = next(
            shape
            for shape in presentation.slides[0].shapes
            if getattr(shape, "text", "") == "感谢聆听"
        )

        assert title_shape.width >= Emu(px_to_emu(4 * 82))


class TestAddToSlide:
    """Tests for adding SVG to existing slides."""

    def test_add_to_existing_slide(self):
        """Test adding SVG shapes to an existing presentation."""
        svg = '''<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        
        # Create a presentation with one slide
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add SVG shapes
        converter = SVGConverter()
        converter.add_string_to_slide(svg, slide)
        
        # Should have at least one shape
        assert len(slide.shapes) >= 1

    def test_add_multiple_svgs(self):
        """Test adding multiple SVGs to the same slide."""
        svg1 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <rect x="0" y="0" width="50" height="50" fill="red"/>
        </svg>'''
        svg2 = '''<svg xmlns="http://www.w3.org/2000/svg" width="50" height="50">
            <circle cx="25" cy="25" r="25" fill="blue"/>
        </svg>'''
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        converter = SVGConverter()
        converter.add_string_to_slide(svg1, slide, x=0, y=0)
        converter.add_string_to_slide(svg2, slide, x=500000, y=0)
        
        # Should have at least 2 shapes
        assert len(slide.shapes) >= 2


class TestConfig:
    """Tests for configuration options."""

    def test_default_config(self):
        config = Config()
        assert config.scale == 1.0
        assert config.preserve_groups is False
        assert config.flatten_groups is True
        assert config.disable_shadows is True

    def test_custom_config(self):
        config = Config(
            scale=2.0,
            curve_tolerance=0.5,
            preserve_groups=False,
        )
        assert config.scale == 2.0
        assert config.curve_tolerance == 0.5
        assert config.preserve_groups is False

    def test_preserve_groups_disables_flatten_groups(self):
        config = Config(preserve_groups=True)

        assert config.preserve_groups is True
        assert config.flatten_groups is False
